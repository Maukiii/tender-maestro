from pathlib import Path


def extract_text(filepath: str) -> str:
    path = Path(filepath)
    ext = path.suffix.lower()

    if ext == ".pdf":
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)

    if ext in (".xlsx", ".xls"):
        if ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"[{sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    lines.append("\t".join("" if v is None else str(v) for v in row))
            return "\n".join(lines)
        else:
            import xlrd
            wb = xlrd.open_workbook(path)
            lines = []
            for sheet in wb.sheets():
                lines.append(f"[{sheet.name}]")
                for i in range(sheet.nrows):
                    lines.append("\t".join(str(v) for v in sheet.row_values(i)))
            return "\n".join(lines)

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)
        return "\n".join(lines)

    if ext == ".docx":
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext in (".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".yaml", ".yml", ".log"):
        return path.read_text(encoding="utf-8", errors="replace")

    raise ValueError(f"Unsupported file type: '{ext}'")
